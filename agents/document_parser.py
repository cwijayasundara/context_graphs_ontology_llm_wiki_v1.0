"""Document parsing adapters for raw/ artifacts.

The DeepAgent ingestor needs document structure, not raw PDF bytes. This module
normalizes local parsing and external document-intelligence providers into one
cacheable JSON/Markdown representation.
"""
from __future__ import annotations

import asyncio
import dataclasses
import json
import os
import re
from pathlib import Path
from typing import Any, Callable, Protocol

import requests

import _lib  # noqa: F401  — triggers .env load

ROOT = Path(__file__).resolve().parents[1]
RAW = ROOT / "raw"
PARSED = ROOT / "parsed"


@dataclasses.dataclass
class ParsedDocument:
    raw_path: str
    provider: str
    markdown: str
    chunks: list[dict[str, Any]]
    metadata: dict[str, Any]
    grounding: dict[str, Any]
    warnings: list[str]

    def to_dict(self) -> dict[str, Any]:
        return dataclasses.asdict(self)


class DocumentParser(Protocol):
    provider: str

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        """Parse a raw artifact and return the canonical representation."""


def _artifact_dir(raw_path: str, parsed_root: Path = PARSED) -> Path:
    rel = Path(raw_path)
    stem = rel.with_suffix("")
    return parsed_root / stem


def write_parsed_artifacts(parsed: ParsedDocument, *, parsed_root: Path = PARSED) -> Path:
    out_dir = _artifact_dir(parsed.raw_path, parsed_root)
    out_dir.mkdir(parents=True, exist_ok=True)
    (out_dir / "document.json").write_text(
        json.dumps(parsed.to_dict(), indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    (out_dir / "document.md").write_text(parsed.markdown.rstrip() + "\n", encoding="utf-8")
    return out_dir


def load_parsed_artifacts(raw_path: str, *, parsed_root: Path = PARSED) -> ParsedDocument | None:
    doc_json = _artifact_dir(raw_path, parsed_root) / "document.json"
    if not doc_json.exists():
        return None
    return ParsedDocument(**json.loads(doc_json.read_text(encoding="utf-8")))


class LiteParseParser:
    provider = "liteparse"

    def __init__(
        self,
        *,
        client: Any | None = None,
        cli_path: str | None = None,
        install_if_not_available: bool | None = None,
        ocr_enabled: bool | None = None,
        ocr_server_url: str | None = None,
        ocr_language: str | None = None,
        dpi: int | None = None,
        target_pages: str | None = None,
        password: str | None = None,
    ) -> None:
        self.client = client
        self.cli_path = cli_path or os.environ.get("LITEPARSE_CLI_PATH")
        self.install_if_not_available = (
            install_if_not_available
            if install_if_not_available is not None
            else os.environ.get("LITEPARSE_INSTALL_IF_MISSING", "false").lower() == "true"
        )
        self.ocr_enabled = _optional_bool(ocr_enabled, "LITEPARSE_OCR_ENABLED")
        self.ocr_server_url = ocr_server_url or os.environ.get("LITEPARSE_OCR_SERVER_URL")
        self.ocr_language = ocr_language or os.environ.get("LITEPARSE_OCR_LANGUAGE")
        self.dpi = dpi or _optional_int("LITEPARSE_DPI")
        self.target_pages = target_pages or os.environ.get("LITEPARSE_TARGET_PAGES")
        self.password = password if password is not None else os.environ.get("DOCUMENT_PARSER_PDF_PASSWORD")

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        suffix = path.suffix.lower()
        warnings: list[str] = []
        metadata: dict[str, Any] = {
            "filename": path.name,
            "extension": suffix,
            "parser": "liteparse",
        }

        if suffix != ".pdf":
            markdown = path.read_text(encoding="utf-8", errors="replace").strip()
            chunk_type = "text" if suffix != ".csv" else "table"
            chunks = []
            if markdown:
                chunks.append(
                    {
                        "id": "liteparse-0",
                        "type": chunk_type,
                        "markdown": markdown,
                        "grounding": {"page": 0, "confidence": None},
                    }
                )
            return ParsedDocument(
                raw_path=raw_path,
                provider=self.provider,
                markdown=markdown,
                chunks=chunks,
                metadata=metadata,
                grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
                warnings=warnings,
            )

        result = self._client().parse(str(path), **self._parse_kwargs())
        markdown = str(_attr_or_key(result, "text", "") or "").strip()
        pages = list(_attr_or_key(result, "pages", []) or [])
        chunks = []
        for idx, page in enumerate(pages):
            page_num = int(_attr_or_key(page, "pageNum", _attr_or_key(page, "page", idx + 1)))
            page_text = str(_attr_or_key(page, "text", "") or "").strip()
            text_items = [_text_item(item) for item in (_attr_or_key(page, "textItems", []) or [])]
            grounding = {
                "page": page_num,
                "width": _attr_or_key(page, "width", None),
                "height": _attr_or_key(page, "height", None),
                "confidence": None,
                "text_items": text_items,
            }
            if page_text:
                chunks.append(
                    {
                        "id": f"liteparse-page-{page_num}",
                        "type": "page",
                        "markdown": page_text,
                        "grounding": grounding,
                    }
                )
        if not markdown and chunks:
            markdown = "\n\n---\n\n".join(chunk["markdown"] for chunk in chunks)
        metadata.update(
            {
                "parser": "liteparse",
                "page_count": len(pages),
                "ocr_enabled": self.ocr_enabled,
                "target_pages": self.target_pages,
            }
        )
        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown=markdown,
            chunks=chunks,
            metadata=metadata,
            grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
            warnings=warnings,
        )

    def _client(self) -> Any:
        if self.client is not None:
            return self.client
        try:
            from liteparse import LiteParse
        except ImportError as exc:
            raise RuntimeError(
                "Install LiteParse for local PDF parsing: "
                "npm install -g @llamaindex/liteparse && pip install liteparse"
            ) from exc
        kwargs: dict[str, Any] = {"install_if_not_available": self.install_if_not_available}
        if self.cli_path:
            kwargs["cli_path"] = self.cli_path
        return LiteParse(**kwargs)

    def _parse_kwargs(self) -> dict[str, Any]:
        kwargs: dict[str, Any] = {}
        if self.ocr_enabled is not None:
            kwargs["ocr_enabled"] = self.ocr_enabled
        if self.ocr_server_url:
            kwargs["ocr_server_url"] = self.ocr_server_url
        if self.ocr_language:
            kwargs["ocr_language"] = self.ocr_language
        if self.dpi:
            kwargs["dpi"] = self.dpi
        if self.target_pages:
            kwargs["target_pages"] = self.target_pages
        if self.password:
            kwargs["password"] = self.password
        return kwargs


class LandingAIParser:
    provider = "landingai"

    def __init__(
        self,
        *,
        api_key: str | None = None,
        endpoint: str | None = None,
        model: str | None = None,
        client: Any | None = None,
        timeout: int | None = None,
    ) -> None:
        self.api_key = api_key or os.environ.get("LANDINGAI_API_KEY")
        if not self.api_key:
            raise RuntimeError("LANDINGAI_API_KEY is required for DOCUMENT_PARSER_PROVIDER=landingai")
        self.endpoint = endpoint or os.environ.get(
            "LANDINGAI_PARSE_URL",
            "https://api.va.landing.ai/v1/ade/parse",
        )
        self.model = model or os.environ.get("LANDINGAI_PARSE_MODEL", "dpt-2-latest")
        self.client = client or requests.Session()
        self.timeout = timeout or int(os.environ.get("DOCUMENT_PARSER_TIMEOUT", "180"))

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        data = {"model": self.model}
        if split := os.environ.get("LANDINGAI_SPLIT"):
            data["split"] = split
        if password := os.environ.get("DOCUMENT_PARSER_PDF_PASSWORD"):
            data["password"] = password

        with path.open("rb") as document:
            response = self.client.post(
                self.endpoint,
                headers={"Authorization": f"Bearer {self.api_key}"},
                files={"document": (path.name, document)},
                data=data,
                timeout=self.timeout,
            )
        response.raise_for_status()
        payload = response.json()
        grounding = payload.get("grounding") or {}
        chunks = []
        for chunk in payload.get("chunks") or []:
            chunk_grounding = dict(chunk.get("grounding") or {})
            chunk_grounding.update(grounding.get(chunk.get("id"), {}))
            chunks.append(
                {
                    "id": chunk.get("id"),
                    "type": chunk.get("type"),
                    "markdown": chunk.get("markdown", ""),
                    "grounding": chunk_grounding,
                }
            )

        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown=(payload.get("markdown") or "").strip(),
            chunks=chunks,
            metadata=payload.get("metadata") or {},
            grounding=grounding,
            warnings=[],
        )


class NemotronParser:
    provider = "nemotron"

    def __init__(
        self,
        *,
        parse_pdf_to_pages: Callable[[Path], Any] | None = None,
    ) -> None:
        self._parse_pdf_to_pages = parse_pdf_to_pages

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        pages = asyncio.run(self._parse_pages(path))
        chunks: list[dict[str, Any]] = []
        markdown_pages: list[str] = []
        for idx, page in enumerate(pages):
            text = _page_text(page).strip()
            if not text:
                continue
            page_number = _page_number(page, idx)
            chunk_id = f"nemotron-page-{page_number}"
            markdown_pages.append(text)
            chunks.append(
                {
                    "id": chunk_id,
                    "type": "page",
                    "markdown": text,
                    "grounding": {"page": page_number, "confidence": None},
                }
            )

        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown="\n\n---\n\n".join(markdown_pages).strip(),
            chunks=chunks,
            metadata={"filename": path.name, "page_count": len(pages), "parser": "nemotron-parse"},
            grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
            warnings=[],
        )

    async def _parse_pages(self, path: Path) -> list[Any]:
        parse_pdf_to_pages = self._parse_pdf_to_pages
        if parse_pdf_to_pages is None:
            try:
                from paperqa_nemotron import parse_pdf_to_pages
            except ImportError as exc:
                raise RuntimeError(
                    "Install paper-qa-nemotron or provide a local Nemotron parser endpoint "
                    "to use DOCUMENT_PARSER_PROVIDER=nemotron"
                ) from exc
        return await parse_pdf_to_pages(path)


def _page_text(page: Any) -> str:
    if isinstance(page, dict):
        return str(page.get("text") or page.get("markdown") or page.get("content") or "")
    return str(
        getattr(page, "text", None)
        or getattr(page, "markdown", None)
        or getattr(page, "content", None)
        or ""
    )


def _page_number(page: Any, fallback: int) -> int:
    if isinstance(page, dict):
        metadata = page.get("metadata") or {}
        return int(page.get("page") or metadata.get("page") or fallback)
    return int(getattr(page, "page", fallback))


def _attr_or_key(obj: Any, key: str, default: Any = None) -> Any:
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)


def _text_item(item: Any) -> dict[str, Any]:
    x = _attr_or_key(item, "x", None)
    y = _attr_or_key(item, "y", None)
    width = _attr_or_key(item, "width", None)
    height = _attr_or_key(item, "height", None)
    return {
        "text": _attr_or_key(item, "text", ""),
        "box": {
            "left": x,
            "top": y,
            "right": None if x is None or width is None else x + width,
            "bottom": None if y is None or height is None else y + height,
            "width": width,
            "height": height,
        },
    }


def _optional_bool(value: bool | None, env_name: str) -> bool | None:
    if value is not None:
        return value
    raw = os.environ.get(env_name)
    if raw is None or raw == "":
        return None
    return raw.lower() in {"1", "true", "yes", "on"}


def _optional_int(env_name: str) -> int | None:
    raw = os.environ.get(env_name)
    if not raw:
        return None
    return int(raw)


class WordParser:
    """Parse .docx files using python-docx. Heading-based chunking."""

    provider = "word"

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        try:
            import docx
        except ImportError as exc:
            raise RuntimeError(
                "Install python-docx for Word parsing: pip install python-docx"
            ) from exc

        document = docx.Document(str(path))
        sections: list[dict[str, Any]] = []
        current = {"heading": "", "level": 0, "lines": []}

        def _flush() -> None:
            if current["lines"] or current["heading"]:
                sections.append({
                    "heading": current["heading"],
                    "level": current["level"],
                    "markdown": _render_section(current["heading"], current["level"], current["lines"]),
                })

        for para in document.paragraphs:
            text = (para.text or "").strip()
            if not text:
                continue
            style = (para.style.name or "").lower()
            if style.startswith("heading"):
                _flush()
                level = _heading_level(style)
                current = {"heading": text, "level": level, "lines": []}
            else:
                current["lines"].append(text)
        _flush()

        for table_idx, table in enumerate(document.tables, start=1):
            rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
            md = _table_to_markdown(rows)
            if md:
                sections.append({
                    "heading": f"Table {table_idx}",
                    "level": 2,
                    "markdown": f"## Table {table_idx}\n\n{md}",
                })

        chunks = [
            {
                "id": f"word-section-{idx + 1}",
                "type": "section",
                "markdown": s["markdown"],
                "grounding": {"section_index": idx + 1, "heading": s["heading"], "level": s["level"]},
            }
            for idx, s in enumerate(sections) if s["markdown"].strip()
        ]
        markdown = "\n\n".join(c["markdown"] for c in chunks).strip()
        if not markdown:
            markdown = path.read_text(errors="replace").strip()
        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown=markdown,
            chunks=chunks,
            metadata={
                "filename": path.name,
                "extension": ".docx",
                "section_count": len(chunks),
                "table_count": len(document.tables),
                "parser": "python-docx",
            },
            grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
            warnings=[],
        )


class PptxParser:
    """Parse .pptx files using python-pptx. One chunk per slide."""

    provider = "pptx"

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "Install python-pptx for PowerPoint parsing: pip install python-pptx"
            ) from exc

        prs = Presentation(str(path))
        chunks: list[dict[str, Any]] = []
        slide_markdowns: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = ""
            body_lines: list[str] = []
            notes = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    is_title = (
                        shape.is_placeholder
                        and shape.placeholder_format.idx == 0
                    )
                    if is_title and not title:
                        title = text
                    else:
                        bullet = "  " * max(0, para.level) + "- "
                        body_lines.append(f"{bullet}{text}")
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()

            heading = f"# Slide {slide_idx}: {title}" if title else f"# Slide {slide_idx}"
            md_parts = [heading]
            if body_lines:
                md_parts.append("\n".join(body_lines))
            if notes:
                md_parts.append(f"\n_Notes:_ {notes}")
            md = "\n\n".join(md_parts)
            slide_markdowns.append(md)
            chunks.append({
                "id": f"pptx-slide-{slide_idx}",
                "type": "slide",
                "markdown": md,
                "grounding": {
                    "slide_number": slide_idx,
                    "title": title,
                    "has_notes": bool(notes),
                },
            })

        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown="\n\n---\n\n".join(slide_markdowns).strip(),
            chunks=chunks,
            metadata={
                "filename": path.name,
                "extension": ".pptx",
                "slide_count": len(chunks),
                "parser": "python-pptx",
            },
            grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
            warnings=[],
        )


class WebParser:
    """Parse a downloaded HTML file using trafilatura (boilerplate removal +
    article extraction). Reads a sibling `.meta.json` produced by the web
    crawler so URL/fetched_at/parent_url provenance survives in the wiki."""

    provider = "web"

    def parse(self, path: Path, *, raw_path: str) -> ParsedDocument:
        try:
            import trafilatura
        except ImportError as exc:
            raise RuntimeError(
                "Install trafilatura for HTML parsing: pip install trafilatura"
            ) from exc

        html = path.read_text(encoding="utf-8", errors="replace")
        meta_path = path.with_suffix(".meta.json")
        meta: dict[str, Any] = {}
        if meta_path.exists():
            try:
                meta = json.loads(meta_path.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                meta = {}

        markdown = (trafilatura.extract(
            html,
            output_format="markdown",
            include_links=True,
            include_tables=True,
            with_metadata=False,
        ) or "").strip()

        title = ""
        try:
            tmeta = trafilatura.extract_metadata(html)
            if tmeta and getattr(tmeta, "title", None):
                title = tmeta.title or ""
        except Exception:
            pass

        warnings: list[str] = []
        if not markdown:
            warnings.append("trafilatura returned empty extraction; falling back to raw text")
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html, "html.parser")
                markdown = soup.get_text("\n", strip=True)
            except ImportError:
                markdown = re.sub(r"<[^>]+>", "", html)

        # Chunk by markdown headings; everything before the first heading is
        # one "preamble" chunk. Keeps the page navigable for retrieval.
        sections = _split_by_headings(markdown)
        chunks: list[dict[str, Any]] = []
        for idx, section in enumerate(sections, start=1):
            chunks.append({
                "id": f"web-section-{idx}",
                "type": "section",
                "markdown": section["markdown"],
                "grounding": {
                    "section_index": idx,
                    "heading": section["heading"],
                    "level": section["level"],
                    "url": meta.get("url"),
                },
            })
        if not chunks and markdown:
            chunks.append({
                "id": "web-section-1",
                "type": "section",
                "markdown": markdown,
                "grounding": {"section_index": 1, "heading": title or "", "level": 1, "url": meta.get("url")},
            })

        metadata = {
            "filename": path.name,
            "extension": ".html",
            "parser": "trafilatura",
            "url": meta.get("url"),
            "fetched_at": meta.get("fetched_at"),
            "parent_url": meta.get("parent_url"),
            "depth": meta.get("depth"),
            "content_hash": meta.get("content_hash"),
            "title": title,
            "section_count": len(chunks),
        }
        return ParsedDocument(
            raw_path=raw_path,
            provider=self.provider,
            markdown=markdown,
            chunks=chunks,
            metadata=metadata,
            grounding={chunk["id"]: chunk["grounding"] for chunk in chunks},
            warnings=warnings,
        )


_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$", re.MULTILINE)


def _split_by_headings(markdown: str) -> list[dict[str, Any]]:
    if not markdown.strip():
        return []
    matches = list(_HEADING_RE.finditer(markdown))
    if not matches:
        return [{"heading": "", "level": 0, "markdown": markdown.strip()}]

    sections: list[dict[str, Any]] = []
    if matches[0].start() > 0:
        preamble = markdown[: matches[0].start()].strip()
        if preamble:
            sections.append({"heading": "", "level": 0, "markdown": preamble})
    for i, m in enumerate(matches):
        end = matches[i + 1].start() if i + 1 < len(matches) else len(markdown)
        body = markdown[m.start():end].strip()
        if body:
            sections.append({"heading": m.group(2), "level": len(m.group(1)), "markdown": body})
    return sections


def _heading_level(style: str) -> int:
    for tok in style.split():
        if tok.isdigit():
            return int(tok)
    return 1


def _render_section(heading: str, level: int, lines: list[str]) -> str:
    parts: list[str] = []
    if heading:
        parts.append(f"{'#' * max(1, min(6, level))} {heading}")
    body = "\n".join(lines).strip()
    if body:
        parts.append(body)
    return "\n\n".join(parts).strip()


def _table_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header = rows[0]
    sep = ["---"] * len(header)
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    for row in rows[1:]:
        padded = row + [""] * max(0, len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return "\n".join(lines)


# Format-specific parsers chosen by file extension. These run regardless of
# DOCUMENT_PARSER_PROVIDER because they handle formats LiteParse/LandingAI/Nemotron
# don't natively cover.
_EXTENSION_PARSERS: dict[str, Callable[[], DocumentParser]] = {
    ".docx": WordParser,
    ".pptx": PptxParser,
    ".html": WebParser,
    ".htm": WebParser,
}


def get_parser(provider: str | None = None,
                path: Path | None = None) -> DocumentParser:
    """Pick a parser. File extension wins for formats with a dedicated parser
    (Word, PowerPoint); otherwise fall back to the configured provider."""
    if path is not None:
        ext = path.suffix.lower()
        factory = _EXTENSION_PARSERS.get(ext)
        if factory is not None:
            return factory()
    name = (provider or os.environ.get("DOCUMENT_PARSER_PROVIDER", "liteparse")).lower()
    if name in {"local", "liteparse"}:
        return LiteParseParser()
    if name == "landingai":
        return LandingAIParser()
    if name == "nemotron":
        return NemotronParser()
    if name == "word":
        return WordParser()
    if name == "pptx":
        return PptxParser()
    if name == "web":
        return WebParser()
    raise ValueError(
        "unknown DOCUMENT_PARSER_PROVIDER {!r}; expected liteparse, landingai, "
        "nemotron, word, pptx, or web".format(name)
    )


def parse_document(
    raw_path: str,
    *,
    provider: str | None = None,
    force: bool = False,
    parsed_root: Path = PARSED,
) -> ParsedDocument:
    full = (ROOT / raw_path).resolve()
    if not str(full).startswith(str(RAW.resolve())):
        raise ValueError("parse_document is restricted to paths under raw/")
    if not full.exists():
        raise FileNotFoundError(raw_path)

    if not force:
        cached = load_parsed_artifacts(raw_path, parsed_root=parsed_root)
        if cached is not None:
            return cached

    parsed = get_parser(provider, path=full).parse(full, raw_path=raw_path)
    write_parsed_artifacts(parsed, parsed_root=parsed_root)
    return parsed
